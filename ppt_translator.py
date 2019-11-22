## ppt_translator:
# A python script for translating a .ppt or .pptx file

## Setup:
# 1) Run file.reg to enable drag-and-drop conversion.
# 2) Install necessary Python 3.8 dependencies (i.e. pptx, googletrans, pandas,).
# 3) Drag and drop .ppt or .pptx file on top of ppt_translator.py.
# 4) Specify target language.
# 5) Open translated document!

## To do:
# preserve page numbers
# access hidden text frames
# give glossary lang labels
# highlight translated content for dubious translations
# if font is larger check if font of next run matches to concatenate
# replace "& \n" and "&\n" with "& "

## Troubleshooting:
# 1) JSONDecodeError: Expecting value: line 1 column 1 (char 0)
# Is usually caused because your IP address has been blocked by google.
# Try changing your IP address with a vpn or wait awhile.
# 2) TypeError: 'NoneType' object is not iterable
# Could be caused by extra spaces or weird characters within the powerpoint.

# -*- coding: utf-8 -*-
#! python3

# Dependencies:
from pptx import Presentation
from pptx.util import Pt
from googletrans import Translator
from googletrans.constants import LANGUAGES
import pandas
import os
import sys
import re

##2nd Level Helper Function
##takes input text
##returns input text translated to tgtLang
def translate_text(input_text, tgtLang):
    print()
    print("1")
    print(input_text)
    # input_text already translated
    if input_text in query_dict:
        output = query_dict[input_text]        

    # input_text in glossary.EN:
    elif input_text in EN_glossary_as_list:
        output = ZH_glossary_as_list[EN_glossary_as_list.index(input_text)]

    # input_text in glossary.ZH:
    elif input_text in ZH_glossary_as_list:
        output = EN_glossary_as_list[ZH_glossary_as_list.index(input_text)]
    
    # input_text is empty
    elif input_text == "":
        output = input_text
        
    # input_text is skippable characters
    elif skip_regex.match(input_text):
        output = input_text
        
    # input_text is spaces
    elif space_regex.match(input_text):
        output = input_text
        
    # detected lang != srcLang -> skip translation if
    elif srcLang != translator.detect(input_text).lang:
        output = input_text

    # google translation is necessary
    else:
        # remove ending spaces
        while space_regex.match(input_text[-1]):
            input_text = input_text[:-1]
        # remove leading spaces
        while space_regex.match(input_text[0]):
            input_text = input_text[1:]
        # remove double spaces
        while re.search("  ", input_text):
            input_text = input_text.Replace("  ", " ")
        print("2")
        print(input_text)
        # check if glossary can be referenced for en->zh-cn translation
        if srcLang == "en" and tgtLang == "zh-cn":
            # iterator to check input for matches with glossary.EN
            i = 0
            length = len(glossary.EN)
            while i < length:
                term_EN = glossary.EN[i]
                term_ZH = glossary.ZH[i]
                # find each occurence of term_EN
                ##add "flags=re.IGNORECASE" to turn on case sensitivity##
                parts = re.split(term_EN, input_text)
                # replace matches w/ term_ZH
                input_text = term_ZH.join(parts)
                i += 1   

        # check if glossary can be referenced for zh-cn->en translation
        elif srcLang == "zh-cn" and tgtLang == "en":
            # iterator to check input for matches with glossary.EN
            i = 0
            length = len(glossary.ZH)
            while i < length:
                term_EN = glossary.EN[i]
                term_ZH = glossary.ZH[i]
                # find each occurence of term_ZH
                parts = input_text.split(term_ZH)
                # replace matches w/ term_ZH
                input_text = term_EN.join(parts) 
                i += 1
        
        print("3")
        print(input_text)
        print(srcLang)

        #TRANSLATE input_text to tgtLang
        output_raw = translator.translate(input_text, src=srcLang, dest=tgtLang)
        output = output_raw.text
        
        print("4")
        print(output)
        print("\n")
 
    # save input_text and output to dictionary
    query_dict[input_text] = output
    return output

##1st Level Helper Function
def translate_ppt(pptfname, tgtLang):
    prs = Presentation(pptfname)
    prs.save(pre + "_" + tgtLang + ext)
    
    for slide in prs.slides:
        #PROGRESS LOGGING
        curr_slide_num = prs.slides.index(slide)
        total_slide_num = len(prs.slides)
        string1 = ("Translating slide [{0}]  \r".format(str(curr_slide_num + 1) +" / "+str(total_slide_num)))
        sys.stdout.write(string1)
        sys.stdout.flush()
        
        for shape in slide.shapes:
            ##if table
            if shape.has_table:
                table = shape.table
                for cell in table.iter_cells():
                    for paragraph in cell.text_frame.paragraphs:
                        paragraphsaved = paragraph

                        #READ
                        read_text = paragraphsaved.text
                        read_font_size = paragraphsaved.font.size
                        is_bold = paragraphsaved.font.bold
                        is_italic = paragraphsaved.font.italic
                        font_color = None
                        read_font = paragraphsaved.font
                        if not paragraphsaved.font.color.type == None:
                            font_color = read_run.font.color.theme_color
                        # TRANSLATE:
                        translated_text = translate_text(read_text, tgtLang)
                        prs.save(pre + "_" + tgtLang + ext)
                        paragraph.clear()
                        
                        # NEWRUN:
                        write_run = paragraph.add_run()
                        write_run.text = translated_text
                        # WRITE:
                        font = write_run.font
                        font.size = read_font_size
                        font.bold = is_bold
                        font.italic = is_italic
                        if not read_run.font.color.type == None:
                            font.color.theme_color = font_color
                #continue
            ##if no text frame                    
            if not shape.has_text_frame:
                continue
            
            #if text frame
            for paragraph in shape.text_frame.paragraphs:
                paragraphRuns = paragraph.runs
                paragraph.clear()
                for read_run in paragraphRuns:
                    
                    #READ
                    read_text = read_run.text
                    read_font_size = read_run.font.size
                    is_bold = read_run.font.bold
                    is_italic = read_run.font.italic
                    font_color = None
                    read_font = read_run.font
                    if not read_run.font.color.type == None:
                        font_color = read_run.font.color.theme_color
                        
                    # TRANSLATE:
                    translated_text = translate_text(read_text, tgtLang)##error prone##
                    prs.save(pre + "_" + tgtLang + ext)
                    
                    # NEWRUN:
                    write_run = paragraph.add_run()
                    write_run.text = translated_text
                    
                    # WRITE:
                    font = write_run.font
                    font.size = read_font_size
                    font.bold = is_bold
                    font.italic = is_italic
                    if not read_run.font.color.type == None:
                        font.color.theme_color = font_color
    prs.save(pre + "_" + tgtLang + ext)

##MAIN##
#prevents command line from exiting upon error
def show_exception_and_exit(exc_type, exc_value, tb):
    import traceback
    traceback.print_exception(exc_type, exc_value, tb)
    input("Press <return> to exit...")
    sys.exit(-1)
sys.excepthook = show_exception_and_exit

# takes argv[1] as fileDir
if len(sys.argv) > 1:
    fileDir = sys.argv[1]
    print(fileDir)
    print()
else:
    raise Exception('To translate, drag and drop a ppt or pptx file onto ppt_translator.py!')

# import glossary.csv
cwd = os.path.dirname(os.path.realpath(__file__))
glossary_abs_dir = os.path.join(cwd, "glossary.csv")

###convert glossary.csv columns to lists for quick lookup
glossary = pandas.read_csv(glossary_abs_dir, names=['ZH','EN'])
ZH_glossary_as_list = glossary.ZH.values.tolist()
EN_glossary_as_list = glossary.EN.values.tolist()
###TODO: allow for multiple columns/languages

# regex for skipping spaces, tabs and empty strings
skip_regex = re.compile(r'''(
    [ 1234567890-_+=!@#$%^&*()[]{}"':;.,/?<>~`—]+
    )''', re.VERBOSE)
space_regex = re.compile(r'''(
    [\s]+
    )''', re.VERBOSE)

# establish translator function
translator = Translator()

#dictionary for searched srcLang tgtLang data
query_dict ={}

#check if file is ppt or pptx  
pre, ext = os.path.splitext(fileDir)
if ext == '.pptx' or ext == '.ppt':
    srcLang = input("\nSource language?\n(press <return> for input detection)\n" )
    tgtLang = input("\nTarget language?\n" )
    srcLang = srcLang.lower()
    tgtLang = tgtLang.lower()
    for k, v in LANGUAGES.items():
        if tgtLang == v:
            tgtLang = k
        if tgtLang == k:
            tgtLang = k
            
    # while loop to make sure tgtLang variable is contained in LANGUAGES      
    while tgtLang not in LANGUAGES:
        print("\n")
        for v in LANGUAGES.values():
            print(v)
        tgtLang = input("\nPlease choose one of the target languages above.\n" )
        tgtLang = tgtLang.lower()
        for k, v in LANGUAGES.items():
            if tgtLang == v:
                tgtLang = k
            if tgtLang == k:
                tgtLang = k   
    for k, v in LANGUAGES.items():
        if srcLang == v:
            srcLang = k
        if srcLang == k:
            srcLang = k
            
    # while loop to make sure tgtLang variable is contained in LANGUAGES      
    while srcLang not in LANGUAGES:
        print("\n")
        for v in LANGUAGES.values():
            print(v)
        srcLang = input("\nPlease choose one of the target languages above.\n" )
        srcLang = srcLang.lower()
        for k, v in LANGUAGES.items():
            if srcLang == v:
                srcLang = k
            if srcLang == k:
                srcLang = k
    #translate ppt
    translate_ppt(fileDir, tgtLang)
